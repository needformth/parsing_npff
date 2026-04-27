import pandas as pd
import numpy as np
import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

class ExcelToPptx:
    # Разметка таблиц по умолчанию (Показатель: 'столбец начала:столбец конца'). Можно менять
    # По умолчанию считается, что в таблице по показателю есть поля: Место, Название НПФ, Показатель (напр. Активы), Доля на рынке
    cols = {
    'Активы': 'A:D',
    'Клиенты': 'E:H',
    'ПН': 'I:L',
    'ЗЛ': 'M:P',
    'ПР': 'Q:T',
    'Участники': 'U:X',
    'ПДС': 'Y:AB',
    'Участники ПДС': 'AC:AF'
    }
    # Информация для индикатора в презентации [название индикатора, его описание]
    indicator_info = {
        'Активы': ['пенсионные активы', 'Активы в млрд рублей / доля рынка в %'],
        'Клиенты': ['клиенты', 'Клиенты в тыс. человек / доля рынка в %'],
        'ПН': ['пенсионные накопления', 'Активы в млрд рублей / доля рынка в %'],
        'ЗЛ': ['застрахованные лица', 'Лица в тыс. человек / доля рынка в %'],
        'ПР': ['пенсионные резервы' , 'Активы в млрд рублей / доля рынка в %'],
        'Участники': ['участники', 'Участники в тыс. человек / доля рынка в %'],
        'ПДС': ['', ''],
        'Участники ПДС': ['', '']
    }
    @staticmethod
    def extract_dataframes(start: int, # № строки начала топа в экселе
                           end: int,  # № строки конца топа в экселе 
                           our_funds_groups: list[str], # список названий наших фондов
                           data: str, # <filepath>.xlsx 
                           sheet_name: str=None, # название листа
                           cols: dict[str, str]=cols # классовый словарь, содержащий разметку таблиц по показателям
                           ) -> dict[str, pd.DataFrame]:
        '''
        Docstring for extract_dataframes
        start: int – номер строки в экселе, где расположены имена полей (№, группа НПФ, Доля и т.д.)
        end: int –  номер, где записи заканчиваются (там где общий итог)
        our_funds: list – список, который содержит названия наших фондов/групп (по умолчанию Эволюция и Будущее)
        data: str – строка, содержащая ссылку на файл
        sheet_name: str – строка, содержащая название листа
        cols: dict[str, str] – словарь, в котором каждому показателю задаются поля в формате {'Активы': 'A:D', ...}
        rtype: dict[str, DataFrame] – возвращается словарь с датафреймами (№, НПФ, Показатель, Доля)
        Метод нужен для того, чтобы вытаскивать данные из экселя в датафреймы с заданной структурой
        '''
        top_funds = {}
        for type_name, cols in cols.items():
            # Вытаскиваем датафрейм каждого показателя
            if sheet_name is not None:
                df = pd.read_excel(data,
                                    header=start - 1,
                                    nrows=end - start, 
                                    usecols=cols,
                                    sheet_name=sheet_name)
            else:
                df = pd.read_excel(data,
                                   header=start - 1,
                                   nrows=end - start, 
                                   usecols=cols)
            
            # Переназначаем поля
            df.columns = ['№', 'НПФ', 'Показатель', 'Доля']

            # 1 убираем пустые значения, 2 убираем тоталы
            df = df.dropna().loc[df['Доля'] != 1]

            # Присваиваем места в соответствии с долей
            df['№'] = df['Доля'].rank(ascending=False).astype(int)
            # Умножаем долю на 100%
            df['Доля'] = np.round(df['Доля'] * 100, 1)

            # Перехватываем возможные ошибки
            if any(fund not in df['НПФ'].values for fund in our_funds_groups):
                raise ValueError('Нет таких фондов')
            elif len(our_funds_groups) > 2:
                raise ValueError('Указано много фондов')
            
            # Формируем топ 3 из чужих фондов
            alien_funds = df.loc[~df['НПФ'].isin(our_funds_groups)] \
                .sort_values(by='Доля', ascending=False) \
                .iloc[:5 - len(our_funds_groups)]
            # Добавялем туда наши
            top_with_ours = pd.concat([
                alien_funds,
                df.loc[df['НПФ'].isin(our_funds_groups)]
            ],
            ignore_index=True) \
                .sort_values(by='Доля', ascending=False)
            top_funds[type_name] = top_with_ours
        
        return top_funds
    
    @staticmethod
    def update_title(slide, # Слайд
                     date_str: str, # Строковая дата
                     is_funds: bool = False) -> None: # Функция ничего не возвращает, она обновляет заголовки, вставляя дату
        '''
        Docstring for update_title
        
        slide: Объект Presentation.slide[i], с которым идет взаимодействие
        date_str: str – строковое представление даты
        is_funds: bool – булево значение, показывающее, какой из двух заголовков мы редактируем (группы/фонды)
        Обновляет заголовок слайда, подставляя дату и текст в зависимости от типа слайда.
        Устанавливает шрифт Segoe UI, размер 18, цвет #1E14A4.
        '''
        title_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.name == "Title":
                title_shape = shape
                break
        if title_shape is None:
            print("Заголовок не найден")
            return

        if is_funds:
            new_text = f"МЕСТА НА РЫНКЕ НПФ ЭВОЛЮЦИЯ И НПФ БУДУЩЕЕ на {date_str}"
        else:
            new_text = f"МЕСТО НА РЫНКЕ ГРУППЫ ФОНДОВ НА {date_str}"

        title_shape.text = new_text
        # Настройка шрифта для всех параграфов и runs
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x1E, 0x14, 0xA4)

    @staticmethod
    def update_indicator_table(slide, # Слайд
                               quadrant_prefix: str, # Один из строковых префиксов 'LU', 'RU', 'LD', 'RD'
                               places: list[int], # Коллекция с местами наших фондов
                               indicator_key: str # Какой индикатор обновляем
                               ) -> None:
        '''
        Docstring for update_indicator_table
        
        slide: Объект Presentation.slide[i], который обновляется
        quadrant_prefix: str – Один из префиксов LU left upper RU right upper LD left down RD right down
        places: list[int] – Список с местами наших фондов
        indicator_key: str – ключ показателя ('Активы', 'Клиенты', ...)
        Обновляет таблицу-индикатор в заданном квадранте.
        '''
        table_name = f"{quadrant_prefix} Indicator"
        target_table = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE and shape.name == table_name:
                target_table = shape.table
                break
        if target_table is None:
            print(f"Таблица {table_name} не найдена")
            return

        # Получаем описание показателя
        ind_desc, desc_line = ExcelToPptx.indicator_info.get(indicator_key, ['', ''])
        if len(places) == 1:
            first_row_text = f"{places[0]} место на пенсионном рынке – {ind_desc}"
        elif len(places) == 2:
            places.sort()
            first_row_text = f"{places[0]} и {places[1]} места на пенсионном рынке – {ind_desc}"
        else:
            raise ValueError('Фондов может быть не более 2 и как минимум 1')
        # Первая строка (ячейка 0,0)
        cell0 = target_table.cell(0, 0)
        cell0.text = first_row_text
        for paragraph in cell0.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(12)
                run.font.bold = False
                run.font.italic = False
                run.font.color.rgb = RGBColor(0, 0, 0)  # чёрный цвет

        # Вторая строка (ячейка 1,0)
        cell1 = target_table.cell(1, 0)
        cell1.text = desc_line
        for paragraph in cell1.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(10.5)
                run.font.italic = True
                run.font.bold = False
    
    @staticmethod
    def fill_slide(slide, # Слайд
                   indicators_data: dict, # Датафрейм для каждого индикатора с местами НПФ по показателям
                   the_date: str, # Дата в строковом представлении
                   is_funds: bool=False, # Маркер выибрающий, какой слайд редактировать, второй или первый
                   our_funds_group: list = None, # Список наших фондов/группы
                   max_bar_ratio: float = 0.95, # Верхняя возможная граница шкалы
                   min_bar_ratio: float = 0.50) -> None: # Нижняя возможная граница шкалы
        '''
        Docstring for fill_slide
        
        slide: Объект Presentation.slide[i], с которым идет взаимодействие
        indicators_data: dict – Датафрейм для каждого индикатора {показатель: DataFrame с колонками '№', 'НПФ', 'Показатель', 'Доля'}
        the_date: str – Дата в строковом представлении
        is_funds: bool – Маркер выибрающий, какой слайд редактировать, второй или первый
        our_funds_group: list – Список наших фондов/групп
        max_bar_ratio: float – Верхняя возможная граница шкалы
        min_bar_ratio: float – Нижнаяя возможная граница шкалы
        Заполняет слайд данными из indicators_data.
        '''

        # Определяем квадранты и соответствующие префиксы групп
        quadrants = [
            ('LU', 'LU Group'),   # левый верх
            ('RU', 'RU Group'),   # правый верх
            ('LD', 'LD Group'),   # левый низ
            ('RD', 'RD Group')    # правый низ
        ]
        indicators = list(indicators_data.keys())
        if len(indicators) != 4:
            raise ValueError(f"Передано {len(indicators)} показателей, нужно 4")
        
        # Обновляем название слайда
        ExcelToPptx.update_title(slide, the_date, is_funds)
        
        # Проходим по каждому квадранту
        for (prefix, group_prefix), indicator in zip(quadrants, indicators):
            # Достаем датафрейм с этим индикатором
            df = indicators_data[indicator].copy()
            # Сортируем по убыванию доли (лидер сверху)
            df = df.sort_values(by='Доля', ascending=False).reset_index(drop=True)
            # Запоминаем места наших НПФ
            places = df.loc[df['НПФ'].isin(our_funds_group), '№'].tolist()
            # Обновляем места и названия 
            ExcelToPptx.update_indicator_table(slide, prefix, places, indicator)
            # Определяем минимальную и максимальную долю для масштабирования ширины
            max_dolya = df['Доля'].max()
            min_dolya = df['Доля'].min()
            # Если все доли равны, избегаем деления на ноль
            if max_dolya == min_dolya:
                # Тогда все бары будут одинаковой ширины (например, 95%)
                def get_width_ratio(dolya):
                    return max_bar_ratio
            else:
                def get_width_ratio(dolya):
                    # Линейная интерполяция: min_dolya -> min_bar_ratio, max_dolya -> max_bar_ratio
                    ratio = min_bar_ratio + (max_bar_ratio - min_bar_ratio) * (dolya - min_dolya) / (max_dolya - min_dolya)
                    return ratio

            # Для каждой позиции от 1 до 5 (порядок на слайде)
            for position in range(1, 6):
                # Получаем строку данных (индекс position-1 в отсортированном df)
                row = df.iloc[position - 1]
                # Имя группы: например, "LU Group 1"
                group_name = f"{group_prefix} {position}"
                # Ищем группу по имени
                target_group = None
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP and shape.name == group_name:
                        target_group = shape
                        break
                if target_group is None:
                    print(f"Не найдена группа: {group_name}")
                    continue

                # Внутри группы ищем текстовое поле "Fund Name" и подгруппу "Group 26"
                fund_name_shape = None
                group26_shape = None # Это наша шкала
                for shape in target_group.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.name == "Fund Name":
                        fund_name_shape = shape
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP and shape.name == "Group 26":
                        group26_shape = shape
                if fund_name_shape is None or group26_shape is None:
                    print(f"В группе {group_name} не найдены Fund Name или Group 26")
                    continue

                # Обновляем текст названия фонда с его местом (№ из данных)
                fund_name_shape.text = f"{row['№']}. {row['НПФ']}"
                # Настройка шрифта
                for paragraph in fund_name_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Segoe UI"
                        run.font.size = Pt(14)
                        if row['НПФ'] in our_funds_group:
                            run.font.bold = True
                        else:
                            run.font.bold = False

                # Внутри Group 26 ищем Outer и Inner Rectangle
                outer_rect = None
                inner_rect = None
                for shape in group26_shape.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if shape.name == "Outer Rectangle":
                            outer_rect = shape
                        elif shape.name == "Inner Rectangle":
                            inner_rect = shape
                if outer_rect is None or inner_rect is None:
                    print(f"В Group 26 группы {group_name} не найдены прямоугольники")
                    continue

                # Меняем цвета для наших фондов
                if row['НПФ'] in our_funds_group:
                    # Inner Rectangle – зелёный (#00B050)
                    inner_rect.fill.solid()
                    inner_rect.fill.fore_color.rgb = RGBColor(0x00, 0xB0, 0x50)
                    # Outer Rectangle – светло-зелёный (#E2F0D9)
                    outer_rect.fill.solid()
                    outer_rect.fill.fore_color.rgb = RGBColor(0xE2, 0xF0, 0xD9)

                # Обновляем ширину Inner Rectangle пропорционально доле
                width_ratio = get_width_ratio(row['Доля'])
                new_width = int(outer_rect.width * width_ratio)
                inner_rect.width = new_width

                # Форматируем текст внутри бара (Показатель / Доля%)
                value = row['Показатель']
                if isinstance(value, float) and value.is_integer():
                    value_str = f"{int(value):,}".replace(",", " ")
                elif isinstance(value, (int, float)):
                    value_str = f"{value:,.1f}".replace(",", " ")
                else:
                    value_str = str(value)
                percent = row['Доля']
                inner_rect.text = f"{value_str} / {percent}%"

                # Настройка шрифта в баре: Segoe UI 14, жирный, белый
                for paragraph in inner_rect.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Segoe UI"
                        run.font.size = Pt(14)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
    @staticmethod
    def extract_date(data: str,
                     start: int,
                     col_with_date: int=2) -> str:
        '''
        Docstring for extract_date
        
        data: str – путь, в котором расположен эксель, из которого будем вытаскивать информацию
        start: int – строка, в которой начинается таблица
        col_with_date: int – номер столбца, в котором есть дата
        Метод извлекает дату из экселя
        '''
        the_date = pd.read_excel(data,
                    header=start - 2,
                    nrows=1 
                    ).columns[col_with_date-1]
        date_str = datetime.datetime.strftime(the_date, '%d.%m.%Y')
        return date_str

    @staticmethod
    def create_presentation(data: str,
                            groups_position: list|tuple[int],
                            funds_position: list|tuple[int],
                            output_path: str="NPF_parsing_excel2pptx/result/output.pptx", 
                            the_date: str=None,
                            sheet_name: str=None,
                            template_path: str="NPF_parsing_excel2pptx/template.pptx",
                            our_group: list=['Группа Р-Р'],
                            our_funds: list=['Эволюция', 'БУДУЩЕЕ'],
                            indicators: list = ['Активы', 'Клиенты', 'ПН', 'ПР'],
                            max_bar_ratio: float = 0.95, 
                            min_bar_ratio: float = 0.50):
        '''
        Docstring for create_presentation
        
        data: str – директория Excel файла, из которой будут извлекаться данные
        groups_position: list|tuple[int] – № строк экселя, в которых ВКЛЮЧИТЕЛЬНО расположены группы [start, end]
        funds_position: list|tuple[int] – № строк экселя, в которых ВКЛЮЧИТЕЛЬНО расположены группы [start, end]
        ===========Необязательные параметры===========
        output: str – путь, в который будет сохранен pptx файл, напр. output.pptx
        the_date: str – дата в строковом представлении, которая отобразится в презентации
        sheet_name: str – название листа, из которого извлекается информация
        template_path: str – путь образца презентации в формате template.pptx
        our_group: list[str] – список наших групп
        our_funds: list[str] – список наших фондов
        indicators: list[str] – список извлекаемых в презентацию показателей (Активы, Клиенты, ПН, ПР, ЗЛ, Участники, ПДС, Участники ПДС)
        max_bar_ratio: float – доля максимального заполнения шкалы
        min_bar_ratio: float – доля минимального заполнения шкалы
        Основной метод: извлекает данные и заполняет слайд.
        '''


        # Извлекаем данные для всех показателей для групп
        top_groups = ExcelToPptx.extract_dataframes(groups_position[0], groups_position[1], our_group, data, sheet_name)

        # Извлекаем данные для всех показателей для фондов
        top_funds = ExcelToPptx.extract_dataframes(funds_position[0], funds_position[1], our_funds, data, sheet_name)

        if the_date is None:
            the_date = ExcelToPptx.extract_date(data, groups_position[0])

        # Формируем словарь только для нужных показателей
        indicators_data_groups = {}
        indicators_data_funds = {}
        for ind in indicators:
            if ind not in top_funds:
                raise ValueError(f"Показатель '{ind}' не найден в данных")
            indicators_data_groups[ind] = top_groups[ind]
            indicators_data_funds[ind] = top_funds[ind]

        # Загружаем презентацию
        prs = Presentation(template_path)

        # Слайд с группами
        groups_slide = prs.slides[0]
        # Заполняем слайд с группами
        ExcelToPptx.fill_slide(groups_slide, indicators_data_groups, the_date, False, our_group, max_bar_ratio, min_bar_ratio)
        # Слайд с фондами
        funds_slide = prs.slides[1]
        # Заполняем слайд с фондами
        ExcelToPptx.fill_slide(funds_slide, indicators_data_funds, the_date, True, our_funds, max_bar_ratio, min_bar_ratio)
        # Сохраняем
        prs.save(output_path)
        print(f"Презентация сохранена: {output_path}")